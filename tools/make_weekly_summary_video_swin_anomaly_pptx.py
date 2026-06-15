from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
IMG = ROOT / "images"
OUT = ROOT / "presentations" / "weekly_summary_video_swin_anomaly_detection_en.pptx"

WIDE_LAYOUT = (13.333, 7.5)

COLORS = {
    "bg": RGBColor(248, 250, 252),
    "ink": RGBColor(17, 24, 39),
    "muted": RGBColor(75, 85, 99),
    "line": RGBColor(217, 119, 6),
    "blue": RGBColor(219, 234, 254),
    "blue_text": RGBColor(29, 78, 216),
    "green": RGBColor(209, 250, 229),
    "green_text": RGBColor(4, 120, 87),
    "amber": RGBColor(254, 243, 199),
    "amber_text": RGBColor(180, 83, 9),
    "sky": RGBColor(224, 242, 254),
    "sky_text": RGBColor(7, 89, 133),
    "rose": RGBColor(255, 228, 230),
    "rose_text": RGBColor(190, 18, 60),
    "gray": RGBColor(229, 231, 235),
}


def set_background(slide, color=COLORS["bg"]):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_title(slide, text: str):
    box = slide.shapes.add_textbox(Inches(0.65), Inches(0.42), Inches(12.0), Inches(0.62))
    frame = box.text_frame
    frame.clear()
    p = frame.paragraphs[0]
    p.text = text
    p.font.name = "Aptos Display"
    p.font.size = Pt(30)
    p.font.bold = True
    p.font.color.rgb = COLORS["ink"]
    return box


def add_subtitle(slide, text: str, x=0.72, y=2.65, w=10.5, h=0.5):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    p = frame.paragraphs[0]
    p.text = text
    p.font.name = "Aptos"
    p.font.size = Pt(20)
    p.font.color.rgb = COLORS["muted"]
    return box


def add_bullets(slide, bullets: list[str], x, y, w, h, size=20):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    for i, bullet in enumerate(bullets):
        p = frame.paragraphs[0] if i == 0 else frame.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.name = "Aptos"
        p.font.size = Pt(size)
        p.font.color.rgb = COLORS["muted"]
        p.space_after = Pt(8)
    return box


def add_text(slide, text: str, x, y, w, h, size=20, bold=False, color=None, align=None):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    p = frame.paragraphs[0]
    p.text = text
    p.font.name = "Aptos Display" if bold else "Aptos"
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color or COLORS["ink"]
    if align is not None:
        p.alignment = align
    return box


def add_card(slide, x, y, w, h, fill_color):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(x),
        Inches(y),
        Inches(w),
        Inches(h),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = fill_color
    shape.adjustments[0] = 0.08
    return shape


def add_rule(slide, x, y, w, color=COLORS["line"]):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(x),
        Inches(y),
        Inches(w),
        Inches(0.05),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.color.rgb = color
    return shape


def add_picture(slide, name: str, x, y, w, h=None):
    path = IMG / name
    if h is None:
        return slide.shapes.add_picture(str(path), Inches(x), Inches(y), width=Inches(w))
    return slide.shapes.add_picture(str(path), Inches(x), Inches(y), width=Inches(w), height=Inches(h))


def blank_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    return slide


def title_slide(prs):
    slide = blank_slide(prs)
    band = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.9))
    band.fill.solid()
    band.fill.fore_color.rgb = COLORS["ink"]
    band.line.color.rgb = COLORS["ink"]
    add_text(slide, "Weekly Progress Summary", 0.72, 1.75, 11.5, 0.8, 38, True)
    add_subtitle(slide, "Video Swin Transformer and video anomaly detection", y=2.75)
    add_rule(slide, 0.72, 3.58, 2.2)
    add_bullets(
        slide,
        [
            "Completed the next step after Swin Transformer: Video Swin Transformer.",
            "Started the anomaly-detection block with weak supervision and MIL.",
            "Reviewed the first methods that connect surveillance videos to anomaly scores.",
        ],
        0.78,
        4.05,
        11.1,
        1.5,
        21,
    )


def build(prs):
    title_slide(prs)

    slide = blank_slide(prs)
    add_title(slide, "Narrative of the Work")
    add_bullets(
        slide,
        [
            "Continued from ViT, TimeSformer, and Swin Transformer.",
            "Studied how shifted-window attention is extended from images to video.",
            "Moved from action-recognition architectures toward anomaly-detection formulations.",
            "Connected model foundations to the research objective: security-oriented video understanding.",
        ],
        0.85,
        1.32,
        11.4,
        4.6,
        24,
    )

    slide = blank_slide(prs)
    add_title(slide, "Why This Step Matters")
    add_bullets(
        slide,
        [
            "Video Swin Transformer closes the current foundations block on modern video architectures.",
            "Anomaly detection changes the problem from recognizing known actions to finding unusual events.",
            "This is closer to real surveillance, where the system may not know every possible threat in advance.",
            "The new material starts the bridge from model study to research design.",
        ],
        0.75,
        1.25,
        5.9,
        4.8,
        21,
    )
    add_card(slide, 7.05, 1.25, 5.3, 3.9, COLORS["gray"])
    add_text(slide, "Key Idea", 7.35, 1.68, 4.5, 0.45, 27, True)
    add_text(
        slide,
        "The focus is shifting from how to classify actions to how to localize suspicious temporal segments in long surveillance videos.",
        7.35,
        2.42,
        4.55,
        1.7,
        22,
        False,
        COLORS["muted"],
    )

    slide = blank_slide(prs)
    add_title(slide, "Video Swin Transformer")
    add_bullets(
        slide,
        [
            "Extends Swin Transformer from images to video clips.",
            "Uses local attention windows across space and time.",
            "Maintains a hierarchical representation as the video features become deeper.",
            "Offers a practical Transformer backbone for video recognition tasks.",
        ],
        0.75,
        1.25,
        5.1,
        4.7,
        20,
    )
    add_picture(slide, "VideoSwin.png", 6.05, 1.25, 5.9, 3.65)

    slide = blank_slide(prs)
    add_title(slide, "Video Swin Block")
    add_picture(slide, "VSwimBlock.png", 1.1, 1.2, 3.3, 4.4)
    add_bullets(
        slide,
        [
            "The architecture keeps the Swin idea of window-based attention.",
            "Shifted windows help information move across neighboring regions.",
            "The important extension is that attention now operates over video tokens.",
            "This makes Video Swin a natural continuation after TimeSformer and Swin.",
        ],
        5.45,
        1.35,
        6.8,
        4.5,
        21,
    )

    slide = blank_slide(prs)
    add_title(slide, "From Action Recognition to Anomaly Detection")
    add_card(slide, 0.85, 1.55, 3.55, 3.45, COLORS["blue"])
    add_card(slide, 4.9, 1.55, 3.55, 3.45, COLORS["green"])
    add_card(slide, 8.95, 1.55, 3.55, 3.45, COLORS["amber"])
    add_text(slide, "Action Recognition", 1.15, 1.9, 2.95, 0.55, 23, True, COLORS["blue_text"])
    add_text(slide, "Classify known actions in short clips.", 1.15, 2.75, 2.95, 1.0, 21, False, COLORS["muted"])
    add_text(slide, "Temporal Modeling", 5.2, 1.9, 2.95, 0.55, 23, True, COLORS["green_text"])
    add_text(slide, "Learn motion and context over time.", 5.2, 2.75, 2.95, 1.0, 21, False, COLORS["muted"])
    add_text(slide, "Anomaly Detection", 9.25, 1.9, 2.95, 0.55, 23, True, COLORS["amber_text"])
    add_text(slide, "Find abnormal moments in long videos.", 9.25, 2.75, 2.95, 1.0, 21, False, COLORS["muted"])
    add_bullets(
        slide,
        [
            "The research direction now moves from model foundations to surveillance-specific learning problems.",
        ],
        1.0,
        5.55,
        11.1,
        0.7,
        20,
    )

    slide = blank_slide(prs)
    add_title(slide, "Core Anomaly-Detection Problem")
    add_bullets(
        slide,
        [
            "Identify whether a video contains abnormal behavior.",
            "Locate the time interval where the abnormal event occurs.",
            "Optionally classify the type of event after detection.",
            "This fits security footage because videos are long and temporal annotations are expensive.",
        ],
        0.75,
        1.25,
        5.35,
        4.7,
        20,
    )
    add_picture(slide, "AnomalyArch.png", 6.25, 1.28, 5.8, 3.55)

    slide = blank_slide(prs)
    add_title(slide, "Weak Supervision and MIL")
    add_bullets(
        slide,
        [
            "Weak supervision uses video-level labels instead of frame-level annotations.",
            "Multiple Instance Learning treats each video as a bag of temporal segments.",
            "An anomalous video is expected to contain at least one abnormal segment.",
            "The model learns to assign higher anomaly scores to the most suspicious segments.",
        ],
        0.85,
        1.25,
        11.4,
        4.9,
        23,
    )

    slide = blank_slide(prs)
    add_title(slide, "Sultani et al. (2018)")
    add_bullets(
        slide,
        [
            "Introduced a weakly supervised formulation for real-world surveillance anomaly detection.",
            "Uses normal and anomalous video bags with MIL ranking loss.",
            "Works with UCF-Crime: real surveillance videos and 13 anomaly categories.",
            "Provides the conceptual base for later weakly supervised VAD methods.",
        ],
        0.75,
        1.25,
        5.25,
        4.7,
        20,
    )
    add_picture(slide, "AnomalyExamples.png", 6.25, 1.3, 5.8, 3.7)

    slide = blank_slide(prs)
    add_title(slide, "RTFM: Robust Temporal Feature Magnitude")
    add_bullets(
        slide,
        [
            "Builds on weakly supervised video anomaly detection.",
            "Uses feature magnitude as stronger evidence for abnormal snippets.",
            "Aims to separate normal and abnormal temporal features more clearly.",
            "Represents a step beyond traditional MIL scoring strategies.",
        ],
        0.75,
        1.25,
        5.25,
        4.7,
        20,
    )
    add_picture(slide, "RTFM_T1.png", 6.35, 1.25, 5.4, 3.75)

    slide = blank_slide(prs)
    add_title(slide, "VADCLIP")
    add_bullets(
        slide,
        [
            "Explores how CLIP can be adapted to weakly supervised video anomaly detection.",
            "Adds temporal modeling so frame-level visual-language features can support video reasoning.",
            "Uses local and global temporal adaptation to capture short and long-range patterns.",
            "Points toward newer anomaly-detection methods that reuse large pretrained models.",
        ],
        0.75,
        1.25,
        5.35,
        4.7,
        20,
    )
    add_picture(slide, "VADCLIP.png", 6.3, 1.35, 5.7, 3.55)

    slide = blank_slide(prs)
    add_title(slide, "Evolution of Ideas")
    add_card(slide, 0.55, 2.05, 2.05, 1.25, COLORS["blue"])
    add_card(slide, 3.0, 2.05, 2.05, 1.25, COLORS["green"])
    add_card(slide, 5.45, 2.05, 2.05, 1.25, COLORS["amber"])
    add_card(slide, 7.9, 2.05, 2.05, 1.25, COLORS["sky"])
    add_card(slide, 10.35, 2.05, 2.45, 1.25, COLORS["rose"])
    add_text(slide, "Swin", 0.75, 2.43, 1.65, 0.4, 20, True, COLORS["blue_text"], PP_ALIGN.CENTER)
    add_text(slide, "Video Swin", 3.12, 2.43, 1.8, 0.4, 19, True, COLORS["green_text"], PP_ALIGN.CENTER)
    add_text(slide, "MIL", 5.65, 2.43, 1.65, 0.4, 20, True, COLORS["amber_text"], PP_ALIGN.CENTER)
    add_text(slide, "RTFM", 8.12, 2.43, 1.65, 0.4, 20, True, COLORS["sky_text"], PP_ALIGN.CENTER)
    add_text(slide, "VADCLIP", 10.62, 2.43, 1.85, 0.4, 20, True, COLORS["rose_text"], PP_ALIGN.CENTER)
    add_bullets(
        slide,
        [
            "From hierarchical attention for images.",
            "To hierarchical attention for videos.",
            "To weakly supervised anomaly scores over temporal segments.",
            "To stronger temporal features and pretrained visual-language representations.",
        ],
        1.05,
        4.1,
        11.3,
        1.8,
        21,
    )

    slide = blank_slide(prs)
    add_title(slide, "What Is Already in Place")
    add_bullets(
        slide,
        [
            "A complete path from CNN-based video recognition to Transformer-based video backbones.",
            "A clear next architecture after Swin: Video Swin Transformer.",
            "A first anomaly-detection formulation based on weak labels and MIL.",
            "Initial coverage of UCF-Crime, RTFM, and VADCLIP.",
            "A stronger bridge between foundations and the surveillance research objective.",
        ],
        0.85,
        1.3,
        11.4,
        4.8,
        22,
    )

    slide = blank_slide(prs)
    add_title(slide, "Suggested Next Step")
    add_bullets(
        slide,
        [
            "Create a compact comparison table for Sultani 2018, RTFM, and VADCLIP.",
            "Compare supervision type, input features, temporal modeling, datasets, and metrics.",
            "Decide which method is the best first implementation target.",
            "Connect the selected method to the broader surveillance pipeline in the roadmap.",
        ],
        0.85,
        1.3,
        11.4,
        4.8,
        22,
    )


def main():
    prs = Presentation()
    prs.slide_width = Inches(WIDE_LAYOUT[0])
    prs.slide_height = Inches(WIDE_LAYOUT[1])
    build(prs)
    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    print(OUT)
    print(f"{len(prs.slides)} slides")


if __name__ == "__main__":
    main()
