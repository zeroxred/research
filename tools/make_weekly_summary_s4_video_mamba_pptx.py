from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
IMG = ROOT / "images"
OUT = ROOT / "presentations" / "05_weekly_summary_s4_to_video_mamba_en.pptx"

WIDE_LAYOUT = (13.333, 7.5)

COLORS = {
    "bg": RGBColor(248, 250, 252),
    "ink": RGBColor(17, 24, 39),
    "muted": RGBColor(75, 85, 99),
    "line": RGBColor(217, 119, 6),
    "blue": RGBColor(219, 234, 254),
    "blue_text": RGBColor(29, 78, 216),
    "green": RGBColor(209, 250, 229),
    "green_text": RGBColor(4, 120, 87),
    "amber": RGBColor(254, 243, 199),
    "amber_text": RGBColor(180, 83, 9),
    "sky": RGBColor(224, 242, 254),
    "sky_text": RGBColor(7, 89, 133),
    "rose": RGBColor(255, 228, 230),
    "rose_text": RGBColor(190, 18, 60),
    "gray": RGBColor(229, 231, 235),
}


def set_background(slide, color=COLORS["bg"]):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_title(slide, text: str):
    box = slide.shapes.add_textbox(Inches(0.65), Inches(0.42), Inches(12.0), Inches(0.62))
    frame = box.text_frame
    frame.clear()
    p = frame.paragraphs[0]
    p.text = text
    p.font.name = "Aptos Display"
    p.font.size = Pt(30)
    p.font.bold = True
    p.font.color.rgb = COLORS["ink"]
    return box


def add_subtitle(slide, text: str, x=0.72, y=2.65, w=10.8, h=0.5):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    p = frame.paragraphs[0]
    p.text = text
    p.font.name = "Aptos"
    p.font.size = Pt(20)
    p.font.color.rgb = COLORS["muted"]
    return box


def add_bullets(slide, bullets: list[str], x, y, w, h, size=20):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    for i, bullet in enumerate(bullets):
        p = frame.paragraphs[0] if i == 0 else frame.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.name = "Aptos"
        p.font.size = Pt(size)
        p.font.color.rgb = COLORS["muted"]
        p.space_after = Pt(8)
    return box


def add_text(slide, text: str, x, y, w, h, size=20, bold=False, color=None, align=None):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    p = frame.paragraphs[0]
    p.text = text
    p.font.name = "Aptos Display" if bold else "Aptos"
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color or COLORS["ink"]
    if align is not None:
        p.alignment = align
    return box


def add_card(slide, x, y, w, h, fill_color):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(x),
        Inches(y),
        Inches(w),
        Inches(h),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = fill_color
    shape.adjustments[0] = 0.08
    return shape


def add_rule(slide, x, y, w, color=COLORS["line"]):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(x),
        Inches(y),
        Inches(w),
        Inches(0.05),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.color.rgb = color
    return shape


def add_picture(slide, name: str, x, y, w, h=None):
    path = IMG / name
    if h is None:
        return slide.shapes.add_picture(str(path), Inches(x), Inches(y), width=Inches(w))
    return slide.shapes.add_picture(str(path), Inches(x), Inches(y), width=Inches(w), height=Inches(h))


def blank_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    return slide


def title_slide(prs):
    slide = blank_slide(prs)
    band = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.9))
    band.fill.solid()
    band.fill.fore_color.rgb = COLORS["ink"]
    band.line.color.rgb = COLORS["ink"]
    add_text(slide, "Weekly Progress Summary", 0.72, 1.75, 11.5, 0.8, 38, True)
    add_subtitle(slide, "From S4 state space models to Video Mamba", y=2.75)
    add_rule(slide, 0.72, 3.58, 2.2)
    add_bullets(
        slide,
        [
            "Started the State Space Model block with S4 as the foundation.",
            "Moved from fixed state-space dynamics to selective SSMs in Mamba.",
            "Connected Mamba to vision and video through Vision Mamba and Video Mamba.",
        ],
        0.78,
        4.05,
        11.1,
        1.5,
        21,
    )


def build(prs):
    title_slide(prs)

    slide = blank_slide(prs)
    add_title(slide, "Narrative of the Work")
    add_bullets(
        slide,
        [
            "Continued from Transformer-based video models into State Space Models.",
            "Reviewed the SSM view: inputs, hidden states, dynamics, and outputs.",
            "Studied how S4 makes long-sequence modeling practical through structured state spaces.",
            "Then moved to Mamba, where the model selectively updates its state based on the input.",
        ],
        0.85,
        1.32,
        11.4,
        4.6,
        24,
    )

    slide = blank_slide(prs)
    add_title(slide, "Why This Step Matters")
    add_bullets(
        slide,
        [
            "Transformers are powerful but attention grows expensive with long sequences.",
            "State Space Models offer a different path: recurrent dynamics with efficient sequence processing.",
            "Mamba keeps the long-context advantage while adding input-dependent selection.",
            "This is relevant for video because long clips require temporal memory without excessive cost.",
        ],
        0.75,
        1.25,
        5.9,
        4.8,
        21,
    )
    add_card(slide, 7.05, 1.25, 5.3, 3.9, COLORS["gray"])
    add_text(slide, "Key Idea", 7.35, 1.68, 4.5, 0.45, 27, True)
    add_text(
        slide,
        "The research path is shifting from attention-based temporal modeling to state-based sequence modeling.",
        7.35,
        2.42,
        4.55,
        1.5,
        23,
        False,
        COLORS["muted"],
    )

    slide = blank_slide(prs)
    add_title(slide, "State Space Models")
    add_bullets(
        slide,
        [
            "An SSM maps an input signal into a hidden state and then projects it to an output.",
            "The hidden state stores information needed to describe the sequence over time.",
            "In discrete form, the model can be interpreted as a recurrent update.",
            "This provides a mathematical foundation for long-range temporal modeling.",
        ],
        0.85,
        1.25,
        11.4,
        3.2,
        22,
    )
    add_card(slide, 1.1, 4.8, 5.0, 1.25, COLORS["blue"])
    add_card(slide, 7.2, 4.8, 5.0, 1.25, COLORS["green"])
    add_text(slide, "Continuous view", 1.35, 5.08, 4.5, 0.3, 22, True, COLORS["blue_text"], PP_ALIGN.CENTER)
    add_text(slide, "x'(t) = Ax(t) + Bu(t)", 1.35, 5.45, 4.5, 0.3, 18, False, COLORS["muted"], PP_ALIGN.CENTER)
    add_text(slide, "Discrete view", 7.45, 5.08, 4.5, 0.3, 22, True, COLORS["green_text"], PP_ALIGN.CENTER)
    add_text(slide, "x_k = A_bar x_{k-1} + B_bar u_k", 7.45, 5.45, 4.5, 0.3, 18, False, COLORS["muted"], PP_ALIGN.CENTER)

    slide = blank_slide(prs)
    add_title(slide, "S4: Structured State Spaces")
    add_bullets(
        slide,
        [
            "S4 is designed to make State Space Models practical for very long sequences.",
            "The same model can be viewed as recurrent for inference and convolutional for parallel training.",
            "This is the important bridge: stateful memory without giving up efficient computation.",
            "S4 establishes the foundation that Mamba later modifies with selection.",
        ],
        0.85,
        1.28,
        11.4,
        4.9,
        23,
    )

    slide = blank_slide(prs)
    add_title(slide, "From S4 to Mamba")
    add_card(slide, 0.85, 1.55, 3.55, 3.45, COLORS["blue"])
    add_card(slide, 4.9, 1.55, 3.55, 3.45, COLORS["green"])
    add_card(slide, 8.95, 1.55, 3.55, 3.45, COLORS["amber"])
    add_text(slide, "S4", 1.15, 1.95, 2.95, 0.45, 25, True, COLORS["blue_text"])
    add_text(slide, "Structured long-sequence memory.", 1.15, 2.72, 2.95, 0.9, 21, False, COLORS["muted"])
    add_text(slide, "Selection", 5.2, 1.95, 2.95, 0.45, 25, True, COLORS["green_text"])
    add_text(slide, "Input-dependent state updates.", 5.2, 2.72, 2.95, 0.9, 21, False, COLORS["muted"])
    add_text(slide, "Mamba", 9.25, 1.95, 2.95, 0.45, 25, True, COLORS["amber_text"])
    add_text(slide, "Linear-time sequence modeling.", 9.25, 2.72, 2.95, 0.9, 21, False, COLORS["muted"])
    add_bullets(
        slide,
        [
            "The central change is that Mamba decides what information to keep or update based on the current token.",
        ],
        1.0,
        5.55,
        11.1,
        0.7,
        20,
    )

    slide = blank_slide(prs)
    add_title(slide, "Mamba: Selective State Spaces")
    add_bullets(
        slide,
        [
            "Mamba introduces selection into State Space Models.",
            "The model can make its dynamics depend on the input sequence.",
            "This helps it filter relevant information and ignore less useful tokens.",
            "The result is a sequence model that competes with attention while scaling linearly.",
        ],
        0.75,
        1.2,
        5.15,
        4.7,
        20,
    )
    add_picture(slide, "SELECTIONMAMBA.png", 6.05, 1.5, 6.0, 1.85)
    add_picture(slide, "MAMBA_SLAWS.png", 6.05, 4.05, 6.0, 1.9)

    slide = blank_slide(prs)
    add_title(slide, "Mamba Block")
    add_picture(slide, "MAMBABLOCK.png", 0.95, 1.18, 4.25, 3.65)
    add_bullets(
        slide,
        [
            "The block combines projection, selective SSM computation, gating, and output projection.",
            "Gating controls how much information passes through the block.",
            "The selective scan allows recurrent-style processing to remain efficient.",
            "This block becomes the reusable unit for later vision and video variants.",
        ],
        6.05,
        1.3,
        5.9,
        4.7,
        21,
    )

    slide = blank_slide(prs)
    add_title(slide, "Vision Mamba")
    add_bullets(
        slide,
        [
            "Vision Mamba adapts Mamba from 1-D sequences to image tokens.",
            "The goal is to preserve Transformer-like representation power with linear complexity.",
            "Bidirectional scanning helps visual tokens exchange information across the image.",
            "This creates the bridge from language sequence modeling to visual representation learning.",
        ],
        0.75,
        1.2,
        5.15,
        4.7,
        20,
    )
    add_picture(slide, "VISIONMAMBA.png", 6.05, 1.45, 6.0, 2.35)

    slide = blank_slide(prs)
    add_title(slide, "VMamba and Visual State-Space Models")
    add_picture(slide, "VMAMBA_FIG1.png", 0.75, 1.25, 6.0, 2.65)
    add_bullets(
        slide,
        [
            "Visual Mamba variants explore scan directions and image-specific sequence ordering.",
            "The main challenge is turning 2-D visual structure into efficient sequences.",
            "This step prepares the same idea for video, where the model must handle space and time.",
        ],
        7.15,
        1.3,
        5.2,
        4.3,
        21,
    )

    slide = blank_slide(prs)
    add_title(slide, "Video Mamba")
    add_bullets(
        slide,
        [
            "Video Mamba extends Mamba-style sequence modeling to video understanding.",
            "Video tokens require both spatial and temporal ordering.",
            "The model uses scan strategies to process video efficiently.",
            "This connects State Space Models directly to the earlier video-recognition roadmap.",
        ],
        0.75,
        1.2,
        5.15,
        4.7,
        20,
    )
    add_picture(slide, "VIDEOMAMBA.png", 6.05, 1.25, 6.0, 2.35)
    add_picture(slide, "VIDEOMAMBA_SCAN.png", 6.05, 4.05, 6.0, 2.1)

    slide = blank_slide(prs)
    add_title(slide, "Video Mamba: Masked Modeling")
    add_bullets(
        slide,
        [
            "Masked modeling trains the model to recover missing visual-temporal information.",
            "This encourages robust video representations without requiring detailed labels.",
            "The idea is useful for learning from large video data where annotation is expensive.",
            "It also connects naturally to future anomaly-detection work.",
        ],
        0.75,
        1.25,
        5.4,
        4.7,
        20,
    )
    add_picture(slide, "VIDEOMAMBA_MASKING.png", 6.25, 1.35, 5.7, 2.65)

    slide = blank_slide(prs)
    add_title(slide, "Evolution of Ideas")
    add_card(slide, 0.55, 2.05, 2.0, 1.25, COLORS["blue"])
    add_card(slide, 2.95, 2.05, 2.0, 1.25, COLORS["green"])
    add_card(slide, 5.35, 2.05, 2.0, 1.25, COLORS["amber"])
    add_card(slide, 7.75, 2.05, 2.0, 1.25, COLORS["sky"])
    add_card(slide, 10.15, 2.05, 2.6, 1.25, COLORS["rose"])
    add_text(slide, "S4", 0.75, 2.43, 1.6, 0.4, 21, True, COLORS["blue_text"], PP_ALIGN.CENTER)
    add_text(slide, "Mamba", 3.15, 2.43, 1.6, 0.4, 21, True, COLORS["green_text"], PP_ALIGN.CENTER)
    add_text(slide, "Vim", 5.55, 2.43, 1.6, 0.4, 21, True, COLORS["amber_text"], PP_ALIGN.CENTER)
    add_text(slide, "VMamba", 7.95, 2.43, 1.6, 0.4, 20, True, COLORS["sky_text"], PP_ALIGN.CENTER)
    add_text(slide, "Video Mamba", 10.34, 2.43, 2.22, 0.4, 20, True, COLORS["rose_text"], PP_ALIGN.CENTER)
    add_bullets(
        slide,
        [
            "From structured state-space memory.",
            "To selective input-dependent sequence modeling.",
            "To image token scanning and visual representations.",
            "To efficient video sequence modeling across space and time.",
        ],
        1.05,
        4.1,
        11.3,
        1.8,
        21,
    )

    slide = blank_slide(prs)
    add_title(slide, "What Is Already in Place")
    add_bullets(
        slide,
        [
            "A conceptual foundation for State Space Models and S4.",
            "A clear transition from fixed SSM dynamics to selective SSMs.",
            "An understanding of the Mamba block and selective scan.",
            "A bridge from Mamba to visual models through Vision Mamba and VMamba.",
            "A first map of how Video Mamba can support long video understanding.",
        ],
        0.85,
        1.3,
        11.4,
        4.8,
        22,
    )

    slide = blank_slide(prs)
    add_title(slide, "Suggested Next Step")
    add_bullets(
        slide,
        [
            "Compare Video Mamba against Transformer-based video models already reviewed.",
            "Focus on temporal modeling, computational complexity, and suitability for long videos.",
            "Decide whether Mamba-based models are useful as backbones for surveillance anomaly detection.",
            "Use this comparison to connect the SSM block back to the main research roadmap.",
        ],
        0.85,
        1.3,
        11.4,
        4.8,
        22,
    )


def main():
    prs = Presentation()
    prs.slide_width = Inches(WIDE_LAYOUT[0])
    prs.slide_height = Inches(WIDE_LAYOUT[1])
    build(prs)
    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    print(OUT)
    print(f"{len(prs.slides)} slides")


if __name__ == "__main__":
    main()
