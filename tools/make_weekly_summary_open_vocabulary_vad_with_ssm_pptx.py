from __future__ import annotations

from pptx import Presentation
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.util import Inches

from make_weekly_summary_open_vocabulary_vad_pptx import (
    COLORS,
    WIDE_LAYOUT,
    add_bullets,
    add_card,
    add_text,
    blank_slide,
    build as build_open_vad,
)

from pathlib import Path


ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "presentations" / "04_weekly_summary_open_vocabulary_vad_with_ssm_en.pptx"


def add_section_slide(prs, title: str, subtitle: str):
    slide = blank_slide(prs)
    band = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.9))
    band.fill.solid()
    band.fill.fore_color.rgb = COLORS["ink"]
    band.line.color.rgb = COLORS["ink"]
    add_text(slide, title, 0.72, 1.8, 11.5, 0.8, 36, True)
    add_text(slide, subtitle, 0.74, 2.7, 10.8, 0.55, 20, False, COLORS["muted"])
    return slide


def append_ssm_section(prs):
    slide = add_section_slide(prs, "State Space Models", "The next block shifts from anomaly semantics to sequence dynamics.")
    add_bullets(
        slide,
        [
            "The question now is how to model long sequences efficiently.",
            "State space models treat state as compressed memory of the past.",
            "Mamba makes the state update input-dependent instead of fixed.",
        ],
        0.82,
        3.75,
        11.2,
        1.6,
        22,
    )

    slide = blank_slide(prs)
    add_title = lambda s, t: add_text(s, t, 0.72, 0.42, 12.0, 0.62, 30, True, COLORS["ink"])
    add_title(slide, "What Is a State?")
    add_bullets(
        slide,
        [
            "A state is the information needed to describe a system at a given moment.",
            "In sequence models, the state acts as compact memory for what has happened so far.",
            "The transition updates that memory as new inputs arrive.",
            "This is the conceptual base before S4 and Mamba.",
        ],
        0.78,
        1.3,
        5.35,
        4.8,
        20,
    )
    add_card(slide, 6.45, 1.5, 5.9, 3.7, COLORS["gray"])
    add_text(slide, "State = memory", 6.9, 1.95, 5.0, 0.45, 28, True)
    add_text(
        slide,
        "The point is not to store everything. The point is to keep only what is useful for predicting the next step.",
        6.9,
        2.7,
        4.95,
        1.5,
        22,
        False,
        COLORS["muted"],
    )

    slide = blank_slide(prs)
    add_title(slide, "S4 vs. Mamba")
    add_card(slide, 0.7, 1.45, 5.8, 3.95, COLORS["blue"])
    add_card(slide, 6.85, 1.45, 5.8, 3.95, COLORS["green"])
    add_text(slide, "S4", 1.05, 1.85, 5.0, 0.45, 26, True, COLORS["blue_text"])
    add_bullets(
        slide,
        [
            "Uses fixed state dynamics.",
            "Efficient for long sequences.",
            "But every token updates memory in the same way.",
        ],
        1.05,
        2.5,
        4.95,
        2.3,
        19,
    )
    add_text(slide, "Mamba", 7.2, 1.85, 5.0, 0.45, 26, True, COLORS["green_text"])
    add_bullets(
        slide,
        [
            "Makes the update input-dependent.",
            "Important tokens can produce larger state changes.",
            "Selective scan keeps the computation efficient.",
        ],
        7.2,
        2.5,
        4.95,
        2.3,
        19,
    )

    slide = blank_slide(prs)
    add_title(slide, "Mamba")
    add_bullets(
        slide,
        [
            "Mamba asks what information deserves to remain in memory.",
            "It introduces selective, time-varying dynamics instead of fixed transitions.",
            "The result is an adaptive hidden state that can emphasize relevant tokens and suppress noise.",
        ],
        0.75,
        1.25,
        11.4,
        2.2,
        22,
    )
    add_card(slide, 0.85, 3.8, 3.55, 1.25, COLORS["amber"])
    add_card(slide, 4.9, 3.8, 3.55, 1.25, COLORS["sky"])
    add_card(slide, 8.95, 3.8, 3.55, 1.25, COLORS["rose"])
    add_text(slide, "Selection", 1.12, 4.14, 2.95, 0.35, 22, True, COLORS["amber_text"],)
    add_text(slide, "Keep useful context", 1.0, 4.55, 3.2, 0.3, 18, False, COLORS["muted"])
    add_text(slide, "Selective scan", 5.2, 4.14, 2.95, 0.35, 22, True, COLORS["sky_text"],)
    add_text(slide, "Efficient execution", 5.08, 4.55, 3.2, 0.3, 18, False, COLORS["muted"])
    add_text(slide, "Adaptive state", 9.25, 4.14, 2.95, 0.35, 22, True, COLORS["rose_text"],)
    add_text(slide, "Input-dependent memory", 9.08, 4.55, 3.2, 0.3, 18, False, COLORS["muted"])

    slide = blank_slide(prs)
    add_title(slide, "Why This Matters for VAD")
    add_bullets(
        slide,
        [
            "Surveillance videos are long, noisy, and temporally uneven.",
            "A state-space backbone could be a cleaner way to model long-range context than a heavy attention stack.",
            "The next experiment is to ask where Mamba should sit in the VAD pipeline: before MIL, inside the temporal encoder, or after CLIP features.",
            "That is the natural next bridge after OV-VAD.",
        ],
        0.78,
        1.28,
        11.4,
        4.9,
        21,
    )


def main():
    prs = Presentation()
    prs.slide_width = Inches(WIDE_LAYOUT[0])
    prs.slide_height = Inches(WIDE_LAYOUT[1])
    build_open_vad(prs)
    append_ssm_section(prs)
    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    print(OUT)
    print(f"{len(prs.slides)} slides")


if __name__ == "__main__":
    main()
