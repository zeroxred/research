from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
IMG = ROOT / "images"
OUT = ROOT / "presentations" / "weekly_summary_video_transformers_en.pptx"

WIDE_LAYOUT = (13.333, 7.5)

COLORS = {
    "bg": RGBColor(248, 250, 252),
    "ink": RGBColor(17, 24, 39),
    "muted": RGBColor(75, 85, 99),
    "line": RGBColor(217, 119, 6),
    "blue": RGBColor(219, 234, 254),
    "blue_text": RGBColor(29, 78, 216),
    "green": RGBColor(209, 250, 229),
    "green_text": RGBColor(4, 120, 87),
    "amber": RGBColor(254, 243, 199),
    "amber_text": RGBColor(180, 83, 9),
    "sky": RGBColor(224, 242, 254),
    "sky_text": RGBColor(7, 89, 133),
    "rose": RGBColor(255, 228, 230),
    "rose_text": RGBColor(190, 18, 60),
    "gray": RGBColor(229, 231, 235),
}


def set_background(slide, color=COLORS["bg"]):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_title(slide, text: str):
    box = slide.shapes.add_textbox(Inches(0.65), Inches(0.42), Inches(12.0), Inches(0.62))
    frame = box.text_frame
    frame.clear()
    p = frame.paragraphs[0]
    p.text = text
    p.font.name = "Aptos Display"
    p.font.size = Pt(30)
    p.font.bold = True
    p.font.color.rgb = COLORS["ink"]
    return box


def add_subtitle(slide, text: str, x=0.72, y=2.65, w=10.5, h=0.5):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    p = frame.paragraphs[0]
    p.text = text
    p.font.name = "Aptos"
    p.font.size = Pt(20)
    p.font.color.rgb = COLORS["muted"]
    return box


def add_bullets(slide, bullets: list[str], x, y, w, h, size=20):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    for i, bullet in enumerate(bullets):
        p = frame.paragraphs[0] if i == 0 else frame.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.name = "Aptos"
        p.font.size = Pt(size)
        p.font.color.rgb = COLORS["muted"]
        p.space_after = Pt(8)
    return box


def add_text(slide, text: str, x, y, w, h, size=20, bold=False, color=None, align=None):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    p = frame.paragraphs[0]
    p.text = text
    p.font.name = "Aptos Display" if bold else "Aptos"
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color or COLORS["ink"]
    if align is not None:
        p.alignment = align
    return box


def add_card(slide, x, y, w, h, fill_color):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(x),
        Inches(y),
        Inches(w),
        Inches(h),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = fill_color
    shape.adjustments[0] = 0.08
    return shape


def add_rule(slide, x, y, w, color=COLORS["line"]):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(x),
        Inches(y),
        Inches(w),
        Inches(0.05),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.color.rgb = color
    return shape


def add_picture(slide, name: str, x, y, w, h=None):
    path = IMG / name
    if h is None:
        return slide.shapes.add_picture(str(path), Inches(x), Inches(y), width=Inches(w))
    return slide.shapes.add_picture(str(path), Inches(x), Inches(y), width=Inches(w), height=Inches(h))


def blank_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    return slide


def title_slide(prs):
    slide = blank_slide(prs)
    band = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.9))
    band.fill.solid()
    band.fill.fore_color.rgb = COLORS["ink"]
    band.line.color.rgb = COLORS["ink"]
    add_text(slide, "Weekly Progress Summary", 0.72, 1.75, 11.5, 0.8, 38, True)
    add_subtitle(slide, "From CNN and ViT foundations to TimeSformer and Swin Transformer", y=2.75)
    add_rule(slide, 0.72, 3.58, 2.2)
    add_bullets(
        slide,
        [
            "Extended the study beyond SlowFast into Transformer-based vision models.",
            "Reviewed how images are converted into patch sequences for ViT.",
            "Connected ViT ideas to video through TimeSformer and hierarchical Swin Transformer designs.",
        ],
        0.78,
        4.05,
        11.1,
        1.5,
        21,
    )


def build(prs):
    title_slide(prs)

    slide = blank_slide(prs)
    add_title(slide, "Narrative of the Work")
    add_bullets(
        slide,
        [
            "Continued from action-recognition foundations: Two-Stream, C3D, and SlowFast.",
            "Reviewed CNN concepts as the baseline language of visual feature extraction.",
            "Moved into Vision Transformers as a different way to represent images.",
            "Studied how Transformer ideas scale from images to video and dense visual tasks.",
        ],
        0.85,
        1.32,
        11.4,
        4.6,
        24,
    )

    slide = blank_slide(prs)
    add_title(slide, "Why This Step Matters")
    add_bullets(
        slide,
        [
            "SlowFast models motion with specialized temporal pathways.",
            "Transformers model relationships through attention instead of only local convolutions.",
            "This opens a path toward long-range spatial and temporal reasoning.",
            "The new notebooks build the bridge from image recognition to video understanding.",
        ],
        0.75,
        1.25,
        5.9,
        4.8,
        21,
    )
    add_card(slide, 7.05, 1.25, 5.3, 3.9, COLORS["gray"])
    add_text(slide, "Key Idea", 7.35, 1.68, 4.5, 0.45, 27, True)
    add_text(
        slide,
        "The focus moved from hand-designed video pathways toward attention-based representations that can connect distant visual evidence.",
        7.35,
        2.42,
        4.55,
        1.7,
        22,
        False,
        COLORS["muted"],
    )

    slide = blank_slide(prs)
    add_title(slide, "CNN Foundations")
    add_bullets(
        slide,
        [
            "Reviewed how convolutional layers detect local visual patterns.",
            "Connected filters, padding, pooling, flattening, and RGB channels.",
            "Clarified why CNNs are strong for images but need extensions for temporal video structure.",
            "This provides the baseline for understanding why ViT changes the representation style.",
        ],
        0.75,
        1.25,
        5.1,
        4.7,
        20,
    )
    add_picture(slide, "CONVRGB.png", 6.1, 1.25, 5.7, 3.6)

    slide = blank_slide(prs)
    add_title(slide, "Vision Transformer (ViT)")
    add_bullets(
        slide,
        [
            "Converts an image into a sequence of fixed-size patches.",
            "Adds positional information and a class token.",
            "Uses a standard Transformer encoder for image classification.",
            "Shows that attention can replace convolution when enough data and compute are available.",
        ],
        0.75,
        1.25,
        5.15,
        4.7,
        20,
    )
    add_picture(slide, "ViT.png", 6.2, 1.35, 5.8, 3.45)

    slide = blank_slide(prs)
    add_title(slide, "From Images to Video Transformers")
    add_card(slide, 0.85, 1.55, 3.55, 3.45, COLORS["blue"])
    add_card(slide, 4.9, 1.55, 3.55, 3.45, COLORS["green"])
    add_card(slide, 8.95, 1.55, 3.55, 3.45, COLORS["amber"])
    add_text(slide, "CNN", 1.15, 1.95, 2.95, 0.45, 25, True, COLORS["blue_text"])
    add_text(slide, "Local visual patterns", 1.15, 2.72, 2.95, 0.85, 21, False, COLORS["muted"])
    add_text(slide, "ViT", 5.2, 1.95, 2.95, 0.45, 25, True, COLORS["green_text"])
    add_text(slide, "Image patches as tokens", 5.2, 2.72, 2.95, 0.85, 21, False, COLORS["muted"])
    add_text(slide, "Video", 9.25, 1.95, 2.95, 0.45, 25, True, COLORS["amber_text"])
    add_text(slide, "Patches across space and time", 9.25, 2.72, 2.95, 1.05, 21, False, COLORS["muted"])
    add_bullets(
        slide,
        [
            "The same token-based view can be extended from image patches to frame-by-frame video patches.",
        ],
        1.0,
        5.55,
        11.1,
        0.7,
        20,
    )

    slide = blank_slide(prs)
    add_title(slide, "TimeSformer")
    add_bullets(
        slide,
        [
            "Adapts the ViT idea directly to video clips.",
            "Represents video as a sequence of patches across frames.",
            "Uses space-time attention to learn visual and temporal relationships.",
            "A key strategy is divided attention: temporal attention followed by spatial attention.",
        ],
        0.8,
        1.25,
        11.4,
        4.9,
        23,
    )

    slide = blank_slide(prs)
    add_title(slide, "TimeSformer: Main Takeaway")
    add_card(slide, 0.9, 1.55, 5.4, 3.6, COLORS["sky"])
    add_card(slide, 7.0, 1.55, 5.4, 3.6, COLORS["rose"])
    add_text(slide, "Temporal Attention", 1.2, 1.95, 4.8, 0.5, 25, True, COLORS["sky_text"])
    add_bullets(slide, ["Looks across frames.", "Captures how a patch changes over time.", "Useful for motion and action dynamics."], 1.2, 2.72, 4.75, 1.9, 19)
    add_text(slide, "Spatial Attention", 7.3, 1.95, 4.8, 0.5, 25, True, COLORS["rose_text"])
    add_bullets(slide, ["Looks within each frame.", "Captures relations between image regions.", "Useful for objects, scene context, and pose."], 7.3, 2.72, 4.75, 1.9, 19)

    slide = blank_slide(prs)
    add_title(slide, "Swin Transformer")
    add_bullets(
        slide,
        [
            "Addresses limitations of standard ViT for dense visual understanding.",
            "Builds hierarchical representations through multiple stages.",
            "Uses local window attention to reduce computational cost.",
            "Uses shifted windows so neighboring windows can exchange information.",
        ],
        0.75,
        1.2,
        5.15,
        4.7,
        20,
    )
    add_picture(slide, "SwinArquitecture.png", 6.1, 1.25, 5.9, 3.75)

    slide = blank_slide(prs)
    add_title(slide, "Swin: Window Attention")
    add_picture(slide, "Swin.png", 0.7, 1.25, 5.6, 3.4)
    add_bullets(
        slide,
        [
            "Attention is computed inside local windows instead of globally.",
            "Shifted windows create cross-window connections in the next layer.",
            "Patch merging gradually reduces resolution and increases feature depth.",
            "This makes the model more practical for larger images and downstream vision tasks.",
        ],
        7.05,
        1.25,
        5.4,
        4.6,
        20,
    )

    slide = blank_slide(prs)
    add_title(slide, "Evolution of Ideas")
    add_card(slide, 0.55, 2.1, 2.25, 1.25, COLORS["blue"])
    add_card(slide, 3.25, 2.1, 2.25, 1.25, COLORS["green"])
    add_card(slide, 5.95, 2.1, 2.25, 1.25, COLORS["amber"])
    add_card(slide, 8.65, 2.1, 2.25, 1.25, COLORS["sky"])
    add_card(slide, 11.35, 2.1, 1.55, 1.25, COLORS["rose"])
    add_text(slide, "CNN", 0.85, 2.48, 1.65, 0.4, 21, True, COLORS["blue_text"], PP_ALIGN.CENTER)
    add_text(slide, "ViT", 3.55, 2.48, 1.65, 0.4, 21, True, COLORS["green_text"], PP_ALIGN.CENTER)
    add_text(slide, "SlowFast", 6.12, 2.48, 1.9, 0.4, 21, True, COLORS["amber_text"], PP_ALIGN.CENTER)
    add_text(slide, "TimeSformer", 8.76, 2.48, 2.0, 0.4, 20, True, COLORS["sky_text"], PP_ALIGN.CENTER)
    add_text(slide, "Swin", 11.48, 2.48, 1.25, 0.4, 20, True, COLORS["rose_text"], PP_ALIGN.CENTER)
    add_bullets(
        slide,
        [
            "From local convolutional features.",
            "To patch-token image representations.",
            "To video models that combine spatial and temporal attention.",
            "To hierarchical attention models that are more efficient and scalable.",
        ],
        1.05,
        4.15,
        11.3,
        1.7,
        21,
    )

    slide = blank_slide(prs)
    add_title(slide, "What Is Already in Place")
    add_bullets(
        slide,
        [
            "A foundation in CNN operations and visual feature extraction.",
            "A clear explanation of ViT as image patches plus Transformer encoder.",
            "A bridge from ViT to video through TimeSformer.",
            "A hierarchical Transformer direction through Swin and shifted windows.",
            "A stronger conceptual base for selecting modern video-recognition architectures.",
        ],
        0.85,
        1.3,
        11.4,
        4.8,
        22,
    )

    slide = blank_slide(prs)
    add_title(slide, "Suggested Next Step")
    add_bullets(
        slide,
        [
            "Connect these architectures back to video anomaly detection.",
            "Compare CNN-based, SlowFast-style, and Transformer-based methods for security footage.",
            "Identify which models are realistic for the available datasets and compute.",
            "Prepare a compact comparison table: input representation, temporal modeling, cost, and expected limitations.",
        ],
        0.85,
        1.3,
        11.4,
        4.8,
        22,
    )


def main():
    prs = Presentation()
    prs.slide_width = Inches(WIDE_LAYOUT[0])
    prs.slide_height = Inches(WIDE_LAYOUT[1])
    build(prs)
    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    print(OUT)
    print(f"{len(prs.slides)} slides")


if __name__ == "__main__":
    main()
