from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
IMG = ROOT / "images"
OUT = ROOT / "presentations" / "weekly_summary_action_recognition_en.pptx"

WIDE_LAYOUT = (13.333, 7.5)

COLORS = {
    "bg": RGBColor(248, 250, 252),
    "ink": RGBColor(17, 24, 39),
    "muted": RGBColor(75, 85, 99),
    "line": RGBColor(217, 119, 6),
    "blue": RGBColor(219, 234, 254),
    "blue_text": RGBColor(29, 78, 216),
    "green": RGBColor(209, 250, 229),
    "green_text": RGBColor(4, 120, 87),
    "amber": RGBColor(254, 243, 199),
    "amber_text": RGBColor(180, 83, 9),
    "sky": RGBColor(224, 242, 254),
    "sky_text": RGBColor(7, 89, 133),
    "gray": RGBColor(229, 231, 235),
}


def set_background(slide, color=COLORS["bg"]):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_title(slide, text: str):
    box = slide.shapes.add_textbox(Inches(0.65), Inches(0.42), Inches(12.0), Inches(0.62))
    frame = box.text_frame
    frame.clear()
    p = frame.paragraphs[0]
    p.text = text
    p.font.name = "Aptos Display"
    p.font.size = Pt(30)
    p.font.bold = True
    p.font.color.rgb = COLORS["ink"]
    return box


def add_subtitle(slide, text: str, x=0.72, y=2.65, w=10.5, h=0.5):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    p = frame.paragraphs[0]
    p.text = text
    p.font.name = "Aptos"
    p.font.size = Pt(20)
    p.font.color.rgb = COLORS["muted"]
    return box


def add_bullets(slide, bullets: list[str], x, y, w, h, size=20):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    for i, bullet in enumerate(bullets):
        p = frame.paragraphs[0] if i == 0 else frame.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.name = "Aptos"
        p.font.size = Pt(size)
        p.font.color.rgb = COLORS["muted"]
        p.space_after = Pt(8)
    return box


def add_text(slide, text: str, x, y, w, h, size=20, bold=False, color=None, align=None):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    p = frame.paragraphs[0]
    p.text = text
    p.font.name = "Aptos Display" if bold else "Aptos"
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color or COLORS["ink"]
    if align is not None:
        p.alignment = align
    return box


def add_card(slide, x, y, w, h, fill_color):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(x),
        Inches(y),
        Inches(w),
        Inches(h),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = fill_color
    shape.adjustments[0] = 0.08
    return shape


def add_rule(slide, x, y, w, color=COLORS["line"]):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(x),
        Inches(y),
        Inches(w),
        Inches(0.05),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.color.rgb = color
    return shape


def add_picture(slide, name: str, x, y, w, h=None):
    path = IMG / name
    if h is None:
        return slide.shapes.add_picture(str(path), Inches(x), Inches(y), width=Inches(w))
    return slide.shapes.add_picture(str(path), Inches(x), Inches(y), width=Inches(w), height=Inches(h))


def blank_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    return slide


def title_slide(prs):
    slide = blank_slide(prs)
    band = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.9))
    band.fill.solid()
    band.fill.fore_color.rgb = COLORS["ink"]
    band.line.color.rgb = COLORS["ink"]
    add_text(slide, "Weekly Progress Summary", 0.72, 1.75, 11.5, 0.8, 38, True)
    add_subtitle(slide, "Action recognition and video anomaly detection for security", y=2.75)
    add_rule(slide, 0.72, 3.58, 2.2)
    add_bullets(
        slide,
        [
            "Reviewed core video-recognition models: Two-Stream, C3D, and SlowFast.",
            "Connected architectural ideas to surveillance and security scenarios.",
            "Mapped the first set of relevant anomaly-detection datasets.",
        ],
        0.78,
        4.05,
        11.1,
        1.5,
        21,
    )


def build(prs):
    title_slide(prs)

    slide = blank_slide(prs)
    add_title(slide, "Narrative of the Work")
    add_bullets(
        slide,
        [
            "Started from CNN foundations and moved toward video understanding.",
            "Studied how models represent motion and temporal structure.",
            "Compared historical approaches: optical flow, 3D convolutions, and multi-rate pathways.",
            "Linked these foundations to surveillance datasets and anomaly detection.",
        ],
        0.85,
        1.32,
        11.4,
        4.6,
        24,
    )

    slide = blank_slide(prs)
    add_title(slide, "Core Problem: Video Is Not Just an Image")
    add_bullets(
        slide,
        [
            "Images mainly provide appearance: objects, people, scene context, and pose.",
            "Videos add motion, temporal order, and causality.",
            "A reversed sequence may represent a different action.",
            "A useful model must capture both spatial and temporal information.",
        ],
        0.75,
        1.25,
        5.9,
        4.8,
        21,
    )
    add_card(slide, 7.05, 1.25, 5.3, 3.9, COLORS["gray"])
    add_text(slide, "Key Idea", 7.35, 1.68, 4.5, 0.45, 27, True)
    add_text(
        slide,
        "In video recognition, what appears in the frame and how it changes over time are both essential.",
        7.35,
        2.42,
        4.55,
        1.5,
        23,
        False,
        COLORS["muted"],
    )

    slide = blank_slide(prs)
    add_title(slide, "Two-Stream Networks (2014)")
    add_bullets(
        slide,
        [
            "Spatial stream: RGB frames for visual appearance.",
            "Temporal stream: optical flow for explicit motion.",
            "Predictions from both streams are fused at the end.",
            "The paper showed that RGB alone is not enough for strong action recognition.",
        ],
        0.75,
        1.25,
        5.0,
        4.7,
        20,
    )
    add_picture(slide, "TwoStream.png", 6.05, 1.35, 6.3, 3.7)

    slide = blank_slide(prs)
    add_title(slide, "Contribution and Limits of Two-Stream Models")
    add_bullets(
        slide,
        [
            "Contribution: clearly separates spatial appearance and temporal motion.",
            "Contribution: made optical flow a central input for deep video models.",
            "Limitation: optical-flow computation is expensive.",
            "Limitation: the pipeline is not fully end-to-end.",
            "Impact: strongly influenced later architectures such as C3D, SlowFast, and video transformers.",
        ],
        0.85,
        1.3,
        11.2,
        4.8,
        22,
    )

    slide = blank_slide(prs)
    add_title(slide, "C3D: 3D Convolutional Networks (2015)")
    add_bullets(
        slide,
        [
            "Extends convolution across height, width, and time.",
            "Learns appearance and motion directly from short video clips.",
            "Does not require precomputed optical flow.",
            "Marks a shift toward end-to-end spatiotemporal learning.",
        ],
        0.75,
        1.2,
        5.15,
        4.7,
        20,
    )
    add_picture(slide, "3DCONV.png", 6.2, 1.35, 5.8, 3.5)

    slide = blank_slide(prs)
    add_title(slide, "C3D Architecture")
    add_picture(slide, "3DCONVARQ.png", 0.7, 1.25, 7.15, 3.85)
    add_bullets(
        slide,
        [
            "8 convolutional layers.",
            "5 max-pooling layers.",
            "3 x 3 x 3 convolution kernels.",
            "Features can transfer to multiple video tasks.",
        ],
        8.05,
        1.25,
        4.6,
        4.6,
        20,
    )

    slide = blank_slide(prs)
    add_title(slide, "Conceptual Comparison")
    add_card(slide, 0.75, 1.4, 5.6, 3.9, COLORS["blue"])
    add_card(slide, 6.95, 1.4, 5.6, 3.9, COLORS["green"])
    add_text(slide, "Two-Stream", 1.05, 1.75, 4.9, 0.5, 27, True, COLORS["blue_text"])
    add_bullets(slide, ["RGB + optical flow", "Motion is computed first", "Late fusion of streams"], 1.05, 2.45, 4.9, 2.3, 20)
    add_text(slide, "C3D", 7.25, 1.75, 4.9, 0.5, 27, True, COLORS["green_text"])
    add_bullets(slide, ["3D convolutions", "Motion is learned directly", "More end-to-end pipeline"], 7.25, 2.45, 4.9, 2.3, 20)

    slide = blank_slide(prs)
    add_title(slide, "SlowFast Networks (2019)")
    add_bullets(
        slide,
        [
            "Starts from a key observation: time is not isotropic like space.",
            "Slow pathway: low frame rate for semantics and appearance.",
            "Fast pathway: high frame rate for fine motion.",
            "Combines both pathways for efficient and accurate action recognition.",
        ],
        0.8,
        1.25,
        11.4,
        4.9,
        23,
    )

    slide = blank_slide(prs)
    add_title(slide, "Evolution of Ideas")
    add_card(slide, 0.9, 2.2, 2.9, 1.3, COLORS["blue"])
    add_card(slide, 5.2, 2.2, 2.9, 1.3, COLORS["green"])
    add_card(slide, 9.5, 2.2, 2.9, 1.3, COLORS["amber"])
    add_text(slide, "Two-Stream", 1.12, 2.58, 2.45, 0.4, 22, True, COLORS["blue_text"], PP_ALIGN.CENTER)
    add_text(slide, "C3D", 5.82, 2.58, 1.9, 0.4, 22, True, COLORS["green_text"], PP_ALIGN.CENTER)
    add_text(slide, "SlowFast", 10.02, 2.58, 1.95, 0.4, 22, True, COLORS["amber_text"], PP_ALIGN.CENTER)
    add_text(slide, "to", 4.02, 2.55, 0.8, 0.4, 24, True, COLORS["muted"], PP_ALIGN.CENTER)
    add_text(slide, "to", 8.32, 2.55, 0.8, 0.4, 24, True, COLORS["muted"], PP_ALIGN.CENTER)
    add_bullets(
        slide,
        [
            "From explicit motion computed with optical flow.",
            "To joint spatiotemporal learning with 3D convolutions.",
            "To architectures that process appearance and motion at different temporal rates.",
        ],
        1.05,
        4.15,
        11.3,
        1.7,
        21,
    )

    slide = blank_slide(prs)
    add_title(slide, "Security-Oriented Datasets Reviewed")
    add_bullets(
        slide,
        [
            "UCF-Crime: primary candidate because it contains real CCTV footage and classes such as robbery and assault.",
            "XD-Violence: strong complementary dataset for violent events, but with movie and YouTube domain bias.",
            "ShanghaiTech: useful academic benchmark for evaluation.",
            "Avenue and UCSD Ped2: useful for quick experiments or pipeline debugging.",
        ],
        0.85,
        1.25,
        11.5,
        5.1,
        21,
    )

    slide = blank_slide(prs)
    add_title(slide, "Initial Dataset Strategy")
    add_card(slide, 0.85, 1.35, 3.7, 3.5, COLORS["green"])
    add_card(slide, 4.9, 1.35, 3.7, 3.5, COLORS["amber"])
    add_card(slide, 8.95, 1.35, 3.7, 3.5, COLORS["sky"])
    add_text(slide, "Primary", 1.15, 1.75, 3.1, 0.4, 23, True, COLORS["green_text"])
    add_text(slide, "UCF-Crime", 1.15, 2.42, 3.0, 1.0, 28, True)
    add_text(slide, "Complementary", 5.2, 1.75, 3.1, 0.4, 22, True, COLORS["amber_text"])
    add_text(slide, "XD-Violence", 5.2, 2.42, 3.0, 1.0, 26, True)
    add_text(slide, "Benchmark", 9.25, 1.75, 3.1, 0.4, 22, True, COLORS["sky_text"])
    add_text(slide, "ShanghaiTech", 9.25, 2.42, 3.0, 1.0, 25, True)
    add_bullets(
        slide,
        [
            "The research contribution will likely come from temporal modeling, architecture, and generalization rather than from the dataset alone.",
        ],
        1.0,
        5.35,
        11.2,
        0.9,
        20,
    )

    slide = blank_slide(prs)
    add_title(slide, "What Is Already in Place")
    add_bullets(
        slide,
        [
            "A conceptual foundation for why video requires temporal modeling.",
            "A historical map of influential action-recognition models.",
            "A comparison between optical flow, 3D CNNs, and SlowFast pathways.",
            "An initial dataset map for surveillance and anomaly-detection research.",
            "A clearer direction toward security-oriented action and anomaly detection.",
        ],
        0.85,
        1.3,
        11.4,
        4.8,
        22,
    )

    slide = blank_slide(prs)
    add_title(slide, "Suggested Next Step")
    add_bullets(
        slide,
        [
            "Review recent methods for video anomaly detection.",
            "Decide whether the project should be supervised, weakly supervised, or self-supervised.",
            "Select the main dataset and evaluation protocol.",
            "Build a comparison table of recent models, metrics, strengths, and limitations.",
        ],
        0.85,
        1.3,
        11.4,
        4.8,
        22,
    )


def main():
    prs = Presentation()
    prs.slide_width = Inches(WIDE_LAYOUT[0])
    prs.slide_height = Inches(WIDE_LAYOUT[1])
    build(prs)
    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    print(OUT)
    print(f"{len(prs.slides)} slides")


if __name__ == "__main__":
    main()
