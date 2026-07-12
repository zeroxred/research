from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
IMG = ROOT / "images"
OUT = ROOT / "presentations" / "04_weekly_summary_open_vocabulary_vad_en.pptx"

WIDE_LAYOUT = (13.333, 7.5)

COLORS = {
    "bg": RGBColor(248, 250, 252),
    "ink": RGBColor(17, 24, 39),
    "muted": RGBColor(75, 85, 99),
    "line": RGBColor(217, 119, 6),
    "blue": RGBColor(219, 234, 254),
    "blue_text": RGBColor(29, 78, 216),
    "green": RGBColor(209, 250, 229),
    "green_text": RGBColor(4, 120, 87),
    "amber": RGBColor(254, 243, 199),
    "amber_text": RGBColor(180, 83, 9),
    "sky": RGBColor(224, 242, 254),
    "sky_text": RGBColor(7, 89, 133),
    "rose": RGBColor(255, 228, 230),
    "rose_text": RGBColor(190, 18, 60),
    "gray": RGBColor(229, 231, 235),
}


def set_background(slide, color=COLORS["bg"]):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_title(slide, text: str):
    box = slide.shapes.add_textbox(Inches(0.65), Inches(0.42), Inches(12.0), Inches(0.62))
    frame = box.text_frame
    frame.clear()
    p = frame.paragraphs[0]
    p.text = text
    p.font.name = "Aptos Display"
    p.font.size = Pt(30)
    p.font.bold = True
    p.font.color.rgb = COLORS["ink"]
    return box


def add_subtitle(slide, text: str, x=0.72, y=2.65, w=10.5, h=0.5):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    p = frame.paragraphs[0]
    p.text = text
    p.font.name = "Aptos"
    p.font.size = Pt(20)
    p.font.color.rgb = COLORS["muted"]
    return box


def add_bullets(slide, bullets: list[str], x, y, w, h, size=20):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    for i, bullet in enumerate(bullets):
        p = frame.paragraphs[0] if i == 0 else frame.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.name = "Aptos"
        p.font.size = Pt(size)
        p.font.color.rgb = COLORS["muted"]
        p.space_after = Pt(8)
    return box


def add_text(slide, text: str, x, y, w, h, size=20, bold=False, color=None, align=None):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    p = frame.paragraphs[0]
    p.text = text
    p.font.name = "Aptos Display" if bold else "Aptos"
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color or COLORS["ink"]
    if align is not None:
        p.alignment = align
    return box


def add_card(slide, x, y, w, h, fill_color):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(x),
        Inches(y),
        Inches(w),
        Inches(h),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = fill_color
    shape.adjustments[0] = 0.08
    return shape


def add_rule(slide, x, y, w, color=COLORS["line"]):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(x),
        Inches(y),
        Inches(w),
        Inches(0.05),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.color.rgb = color
    return shape


def add_picture(slide, name: str, x, y, w, h=None):
    path = IMG / name
    if h is None:
        return slide.shapes.add_picture(str(path), Inches(x), Inches(y), width=Inches(w))
    return slide.shapes.add_picture(str(path), Inches(x), Inches(y), width=Inches(w), height=Inches(h))


def blank_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    return slide


def title_slide(prs):
    slide = blank_slide(prs)
    band = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.9))
    band.fill.solid()
    band.fill.fore_color.rgb = COLORS["ink"]
    band.line.color.rgb = COLORS["ink"]
    add_text(slide, "Weekly Progress Summary", 0.72, 1.75, 11.5, 0.8, 38, True)
    add_subtitle(slide, "From VADCLIP to open-vocabulary video anomaly detection", y=2.75)
    add_rule(slide, 0.72, 3.58, 2.2)
    add_bullets(
        slide,
        [
            "Extended the anomaly-detection block from weak supervision to open-vocabulary reasoning.",
            "Reviewed the gap between closed-set, open-set, and open-vocabulary VAD.",
            "Started organizing how CLIP can support both detection and semantic categorization.",
        ],
        0.78,
        4.05,
        11.1,
        1.5,
        21,
    )


def build(prs):
    title_slide(prs)

    slide = blank_slide(prs)
    add_title(slide, "Narrative of the Work")
    add_bullets(
        slide,
        [
            "Continued from the previous block on Video Swin Transformer, RTFM, and VadCLIP.",
            "Shifted attention from weakly supervised anomaly scoring toward semantic anomaly naming.",
            "Studied open-set and open-vocabulary formulations as the next research step.",
            "Connected the literature to a more realistic surveillance setting where new anomalies can appear.",
        ],
        0.85,
        1.32,
        11.4,
        4.6,
        24,
    )

    slide = blank_slide(prs)
    add_title(slide, "Why This Step Matters")
    add_bullets(
        slide,
        [
            "WSVAD can localize suspicious segments, but it usually stays inside the anomaly classes seen during training.",
            "Open-set VAD improves generalization to unseen anomalies, but still leaves the anomaly type unnamed.",
            "Open-vocabulary VAD adds semantic categorization on top of detection.",
            "That makes the problem closer to how a real monitoring system should behave.",
        ],
        0.75,
        1.25,
        5.9,
        4.8,
        21,
    )
    add_card(slide, 7.05, 1.25, 5.3, 3.9, COLORS["gray"])
    add_text(slide, "Key Idea", 7.35, 1.68, 4.5, 0.45, 27, True)
    add_text(
        slide,
        "The research target moved from asking only whether a video is abnormal to also asking what anomaly it is.",
        7.35,
        2.42,
        4.55,
        1.7,
        22,
        False,
        COLORS["muted"],
    )

    slide = blank_slide(prs)
    add_title(slide, "From Closed Set to Open Vocabulary")
    add_card(slide, 0.7, 1.6, 2.75, 1.35, COLORS["blue"])
    add_card(slide, 3.85, 1.6, 2.75, 1.35, COLORS["green"])
    add_card(slide, 7.0, 1.6, 2.75, 1.35, COLORS["amber"])
    add_card(slide, 10.15, 1.6, 2.45, 1.35, COLORS["rose"])
    add_text(slide, "WSVAD", 1.0, 1.98, 2.1, 0.35, 22, True, COLORS["blue_text"], PP_ALIGN.CENTER)
    add_text(slide, "Detects known anomalies", 0.92, 2.48, 2.25, 0.45, 19, False, COLORS["muted"], PP_ALIGN.CENTER)
    add_text(slide, "Open-Set", 4.12, 1.98, 2.1, 0.35, 22, True, COLORS["green_text"], PP_ALIGN.CENTER)
    add_text(slide, "Detects unseen anomalies", 4.0, 2.48, 2.35, 0.45, 19, False, COLORS["muted"], PP_ALIGN.CENTER)
    add_text(slide, "OV-VAD", 7.28, 1.98, 2.1, 0.35, 22, True, COLORS["amber_text"], PP_ALIGN.CENTER)
    add_text(slide, "Detects and names them", 7.14, 2.48, 2.4, 0.45, 19, False, COLORS["muted"], PP_ALIGN.CENTER)
    add_text(slide, "CLIP", 10.35, 1.98, 2.0, 0.35, 22, True, COLORS["rose_text"], PP_ALIGN.CENTER)
    add_text(slide, "Semantic bridge", 10.32, 2.48, 2.0, 0.45, 19, False, COLORS["muted"], PP_ALIGN.CENTER)
    add_bullets(
        slide,
        [
            "The progression is not only about better anomaly scores, but about richer supervision and richer outputs.",
        ],
        1.0,
        4.35,
        11.2,
        0.7,
        20,
    )

    slide = blank_slide(prs)
    add_title(slide, "VadCLIP")
    add_bullets(
        slide,
        [
            "VadCLIP adapts CLIP to weakly supervised video anomaly detection.",
            "It keeps the binary detection task, but adds semantic alignment between video and text.",
            "This makes it a bridge between MIL-based WSVAD and open-vocabulary anomaly understanding.",
        ],
        0.75,
        1.25,
        5.1,
        4.7,
        20,
    )
    add_picture(slide, "VADCLIP.png", 6.05, 1.28, 6.0, 3.75)

    slide = blank_slide(prs)
    add_title(slide, "VadCLIP: Main Idea")
    add_card(slide, 0.75, 1.45, 5.9, 3.95, COLORS["blue"])
    add_card(slide, 6.7, 1.45, 5.9, 3.95, COLORS["green"])
    add_text(slide, "Detection Branch", 1.1, 1.85, 5.2, 0.45, 26, True, COLORS["blue_text"])
    add_bullets(
        slide,
        [
            "Uses a local-global temporal adapter to model video dynamics.",
            "Produces anomaly confidence scores from frame-level visual features.",
            "Keeps weak supervision through video-level labels and MIL-style reasoning.",
        ],
        1.1,
        2.55,
        5.0,
        2.35,
        19,
    )
    add_text(slide, "Alignment Branch", 7.05, 1.85, 5.2, 0.45, 26, True, COLORS["green_text"])
    add_bullets(
        slide,
        [
            "Uses CLIP text embeddings for anomaly classes.",
            "Adds learnable prompts and visual prompts to refine the class representation.",
            "Aligns selected video snippets with semantic labels through MIL-Align.",
        ],
        7.05,
        2.55,
        5.0,
        2.35,
        19,
    )

    slide = blank_slide(prs)
    add_title(slide, "Open-Vocabulary VAD")
    add_bullets(
        slide,
        [
            "OV-VAD aims to detect anomalies and assign semantic labels to both seen and unseen categories.",
            "The model is built on CLIP to reuse its image-text alignment capability.",
            "The detection branch produces frame-level anomaly confidence scores.",
            "The categorization branch matches video features with text embeddings.",
        ],
        0.75,
        1.25,
        5.15,
        4.7,
        20,
    )
    add_picture(slide, "OVVAD_ARCH.png", 6.05, 1.28, 6.0, 3.75)

    slide = blank_slide(prs)
    add_title(slide, "Overall Framework")
    add_card(slide, 0.75, 1.45, 5.9, 3.95, COLORS["sky"])
    add_card(slide, 6.7, 1.45, 5.9, 3.95, COLORS["green"])
    add_text(slide, "Detection Branch", 1.1, 1.85, 5.2, 0.45, 26, True, COLORS["sky_text"])
    add_bullets(
        slide,
        [
            "Video frames go through CLIP's image encoder.",
            "Temporal Adapter injects video dynamics.",
            "Semantic Knowledge Injection enriches visual features.",
            "A binary detector outputs frame-level anomaly scores.",
        ],
        1.1,
        2.55,
        5.0,
        2.4,
        19,
    )
    add_text(slide, "Categorization Branch", 7.05, 1.85, 5.2, 0.45, 26, True, COLORS["green_text"])
    add_bullets(
        slide,
        [
            "Video features are aggregated into a compact representation.",
            "Category names are encoded with CLIP's text encoder.",
            "Cross-modal alignment predicts the anomaly category.",
            "This is the semantic layer that closed-set WSVAD lacks.",
        ],
        7.05,
        2.55,
        5.0,
        2.4,
        19,
    )

    slide = blank_slide(prs)
    add_title(slide, "Temporal Adapter")
    add_bullets(
        slide,
        [
            "The Temporal Adapter is designed to model frame-to-frame relationships with minimal extra parameters.",
            "It preserves CLIP's pretrained knowledge while adding temporal structure.",
            "The goal is to capture local temporal context without overfitting the seen anomaly classes.",
        ],
        0.75,
        1.25,
        5.2,
        4.7,
        20,
    )
    add_card(slide, 6.35, 1.35, 5.9, 3.8, COLORS["gray"])
    add_text(slide, "Main Role", 6.75, 1.78, 4.9, 0.45, 26, True)
    add_text(
        slide,
        "Bridge image features into a video-aware representation with a light temporal module rather than a heavy video transformer.",
        6.75,
        2.5,
        4.95,
        1.75,
        22,
        False,
        COLORS["muted"],
    )

    slide = blank_slide(prs)
    add_title(slide, "Semantic Knowledge Injection")
    add_bullets(
        slide,
        [
            "SKI injects anomaly-related semantic knowledge from text into the visual stream.",
            "This helps the model reason about what kind of event a frame is related to.",
            "It moves the method beyond pure binary detection.",
        ],
        0.75,
        1.25,
        5.1,
        4.7,
        20,
    )
    add_picture(slide, "OVVAD_T1.png", 6.15, 1.3, 5.6, 3.6)

    slide = blank_slide(prs)
    add_title(slide, "Novel Anomaly Synthesis")
    add_bullets(
        slide,
        [
            "NAS generates synthetic anomalies to help the model handle novel classes.",
            "LLM-guided descriptions are turned into visual content and inserted into normal videos.",
            "This gives the model a way to see categories that were not present in the original training set.",
        ],
        0.75,
        1.25,
        11.4,
        4.9,
        22,
    )

    slide = blank_slide(prs)
    add_title(slide, "Results and Comparison")
    add_bullets(
        slide,
        [
            "The reported results show that the proposed method is competitive with closed-set baselines.",
            "The gain is most relevant on novel anomaly categories, which is the main point of the open-vocabulary setting.",
        ],
        0.75,
        1.2,
        11.5,
        1.35,
        20,
    )
    add_picture(slide, "OVVAD_TABLE1.png", 0.8, 2.35, 11.9, 3.6)

    slide = blank_slide(prs)
    add_title(slide, "Evolution of Ideas")
    add_card(slide, 0.6, 2.05, 2.25, 1.25, COLORS["blue"])
    add_card(slide, 3.3, 2.05, 2.25, 1.25, COLORS["green"])
    add_card(slide, 6.0, 2.05, 2.25, 1.25, COLORS["amber"])
    add_card(slide, 8.7, 2.05, 2.25, 1.25, COLORS["sky"])
    add_card(slide, 11.4, 2.05, 1.45, 1.25, COLORS["rose"])
    add_text(slide, "WSVAD", 0.85, 2.43, 1.7, 0.4, 20, True, COLORS["blue_text"], PP_ALIGN.CENTER)
    add_text(slide, "VadCLIP", 3.5, 2.43, 1.75, 0.4, 20, True, COLORS["green_text"], PP_ALIGN.CENTER)
    add_text(slide, "Open-Set", 6.2, 2.43, 1.75, 0.4, 20, True, COLORS["amber_text"], PP_ALIGN.CENTER)
    add_text(slide, "OV-VAD", 8.85, 2.43, 1.75, 0.4, 20, True, COLORS["sky_text"], PP_ALIGN.CENTER)
    add_text(slide, "CLIP", 11.55, 2.43, 1.1, 0.4, 20, True, COLORS["rose_text"], PP_ALIGN.CENTER)
    add_bullets(
        slide,
        [
            "From weakly supervised anomaly localization.",
            "To CLIP-based visual-language anomaly modeling.",
            "To open-set generalization.",
            "To open-vocabulary detection and semantic categorization.",
        ],
        1.0,
        4.1,
        11.3,
        1.8,
        21,
    )

    slide = blank_slide(prs)
    add_title(slide, "What Is Already in Place")
    add_bullets(
        slide,
        [
            "A coherent path from action recognition to anomaly detection.",
            "A weakly supervised baseline centered on MIL-based reasoning.",
            "A CLIP-based method that already connects visual and textual representations.",
            "A clear next step toward open-vocabulary anomaly understanding.",
        ],
        0.85,
        1.3,
        11.4,
        4.8,
        22,
    )

    slide = blank_slide(prs)
    add_title(slide, "Suggested Next Step")
    add_bullets(
        slide,
        [
            "Move to state space models and Mamba as the next temporal modeling block.",
            "Compare Mamba with the CLIP-based VAD line to see where it can improve sequence efficiency.",
            "Identify how a state-space backbone could fit into weakly supervised or open-vocabulary VAD.",
            "Prepare the next weekly summary around this transition instead of staying only on OV-VAD.",
        ],
        0.85,
        1.3,
        11.4,
        4.8,
        22,
    )


def main():
    prs = Presentation()
    prs.slide_width = Inches(WIDE_LAYOUT[0])
    prs.slide_height = Inches(WIDE_LAYOUT[1])
    build(prs)
    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    print(OUT)
    print(f"{len(prs.slides)} slides")


if __name__ == "__main__":
    main()
